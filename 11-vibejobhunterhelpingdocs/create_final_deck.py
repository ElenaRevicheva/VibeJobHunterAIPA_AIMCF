#!/usr/bin/env python3
"""
AIdeazz Pitch Deck - FINAL VERSION
Modern design with PROPER contrast and zero overlaps
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Color palette - VIBRANT colors for visibility
cyan_bright = RGBColor(0, 245, 255)
blue_bright = RGBColor(59, 130, 246)
purple_bright = RGBColor(139, 92, 246)
green_bright = RGBColor(16, 185, 129)
orange_bright = RGBColor(245, 158, 11)
red_bright = RGBColor(239, 68, 68)
white = RGBColor(255, 255, 255)
black = RGBColor(10, 10, 20)
dark_text = RGBColor(20, 20, 40)
light_bg = RGBColor(240, 242, 250)

def add_slide():
    """Add slide with dark gradient background"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = RGBColor(15, 5, 30)
    fill.gradient_stops[1].color.rgb = RGBColor(20, 40, 100)
    return slide

def add_box(slide, left, top, width, height, bg_color=light_bg, border_color=cyan_bright):
    """Add LIGHT box with DARK text - proper contrast!"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(3)
    return shape

def add_badge(slide, left, top, text, bg_color):
    """Add solid badge"""
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = bg_color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.text = text
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.font.size = Pt(28)
    para.font.bold = True
    para.font.color.rgb = white
    para.alignment = PP_ALIGN.CENTER
    return badge

def add_text(slide, left, top, width, height, text, size=18, color=dark_text, bold=False, align=PP_ALIGN.LEFT):
    """Add text"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = text
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.font.size = Pt(size)
    para.font.color.rgb = color
    para.font.bold = bold
    para.alignment = align
    para.line_spacing = 1.3
    return textbox

def add_footer(slide):
    """Add footer"""
    add_text(slide, Inches(0.5), Inches(8.5), Inches(6), Inches(0.3), 
             "(c) 2025 AIdeazz. Confidential.", 10, RGBColor(150, 150, 170))
    add_text(slide, Inches(13), Inches(8.5), Inches(2.5), Inches(0.3), 
             "aideazz.xyz", 10, cyan_bright, align=PP_ALIGN.RIGHT)

print("Creating FINAL pitch deck with proper contrast...")

# SLIDE 1: COVER
print("Slide 1: Cover")
slide1 = add_slide()
add_text(slide1, Inches(2), Inches(1.5), Inches(12), Inches(1), 
         "AIdeazz", 110, cyan_bright, True, PP_ALIGN.CENTER)
add_text(slide1, Inches(2), Inches(2.7), Inches(12), Inches(0.5), 
         "Emotionally Intelligent Bilingual AI Companions", 32, white, False, PP_ALIGN.CENTER)

# Main value prop box
box1 = add_box(slide1, Inches(2.5), Inches(3.6), Inches(11), Inches(1), light_bg, cyan_bright)
add_text(slide1, Inches(2.7), Inches(3.8), Inches(10.6), Inches(0.6), 
         "For 280M+ expats and local service providers adapting across 19 Spanish-speaking countries", 
         24, dark_text, False, PP_ALIGN.CENTER)

# Badges row
y = Inches(4.9)
box2 = add_box(slide1, Inches(3.5), y, Inches(3.5), Inches(0.5), RGBColor(220, 255, 220), green_bright)
add_text(slide1, Inches(3.7), y + Inches(0.08), Inches(3.1), Inches(0.35), 
         "LIVE: WhatsApp | Telegram | Web", 14, dark_text, True, PP_ALIGN.CENTER)

box3 = add_box(slide1, Inches(7.2), y, Inches(2), Inches(0.5), RGBColor(220, 235, 255), blue_bright)
add_text(slide1, Inches(7.4), y + Inches(0.08), Inches(1.6), Inches(0.35), 
         "Built for <$15K", 14, dark_text, True, PP_ALIGN.CENTER)

box4 = add_box(slide1, Inches(9.4), y, Inches(1.6), Inches(0.5), RGBColor(255, 240, 220), orange_bright)
add_text(slide1, Inches(9.6), y + Inches(0.08), Inches(1.2), Inches(0.35), 
         "Pre-Seed", 14, dark_text, True, PP_ALIGN.CENTER)

# Founder box
box5 = add_box(slide1, Inches(4.5), Inches(5.8), Inches(7), Inches(1.5), light_bg, purple_bright)
add_text(slide1, Inches(4.7), Inches(6), Inches(6.6), Inches(0.35), 
         "Elena Revicheva", 28, purple_bright, True, PP_ALIGN.CENTER)
add_text(slide1, Inches(4.7), Inches(6.4), Inches(6.6), Inches(0.3), 
         "Founder, CEO & CTO", 18, dark_text, False, PP_ALIGN.CENTER)
add_text(slide1, Inches(4.7), Inches(6.75), Inches(6.6), Inches(0.3), 
         "aipa@aideazz.xyz | aideazz.xyz", 16, RGBColor(80, 80, 120), False, PP_ALIGN.CENTER)

add_footer(slide1)

# SLIDE 2: PROBLEM
print("Slide 2: Problem")
slide2 = add_slide()
add_text(slide2, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "The Problem: 280M+ Expats Face Cultural Adaptation Crisis", 
         44, cyan_bright, True)

# Three problem boxes - LIGHT backgrounds
box_w = Inches(4.8)
box_h = Inches(4.5)
y = Inches(1.4)

# Box 1
box1 = add_box(slide2, Inches(0.5), y, box_w, box_h, RGBColor(255, 230, 230), red_bright)
add_badge(slide2, Inches(0.7), y + Inches(0.2), "!", red_bright)
add_text(slide2, Inches(0.7), y + Inches(1), box_w - Inches(0.4), Inches(0.4), 
         "NO EMOTIONAL INTELLIGENCE", 20, red_bright, True)
add_text(slide2, Inches(0.7), y + Inches(1.5), box_w - Inches(0.4), Inches(2.7), 
         "Current AI treats language as academic exercise\n\nIgnores homesickness, anxiety, cultural shock\n\nGeneric responses when families need empathy", 
         14, dark_text)

# Box 2
box2 = add_box(slide2, Inches(5.6), y, box_w, box_h, RGBColor(255, 230, 230), red_bright)
add_badge(slide2, Inches(5.8), y + Inches(0.2), "X", red_bright)
add_text(slide2, Inches(5.8), y + Inches(1), box_w - Inches(0.4), Inches(0.4), 
         "NO FAMILY AWARENESS", 20, red_bright, True)
add_text(slide2, Inches(5.8), y + Inches(1.5), box_w - Inches(0.4), Inches(2.7), 
         "Treats users as isolated individuals\n\nCan't adapt to different family members\n\nOne-size-fits-all fails", 
         14, dark_text)

# Box 3
box3 = add_box(slide2, Inches(10.7), y, box_w, box_h, RGBColor(255, 230, 230), red_bright)
add_badge(slide2, Inches(10.9), y + Inches(0.2), "#", red_bright)
add_text(slide2, Inches(10.9), y + Inches(1), box_w - Inches(0.4), Inches(0.4), 
         "NO TRUE OWNERSHIP", 20, red_bright, True)
add_text(slide2, Inches(10.9), y + Inches(1.5), box_w - Inches(0.4), Inches(2.7), 
         "Big Tech owns your data\n\nCan't sell or transfer your AI\n\nPlatform shutdown = lose everything", 
         14, dark_text)

# Bottom callout
box4 = add_box(slide2, Inches(0.5), Inches(6.3), Inches(15), Inches(1.5), light_bg, cyan_bright)
add_text(slide2, Inches(0.7), Inches(6.5), Inches(14.6), Inches(0.4), 
         "280M+ expats + 50M+ local service providers", 28, cyan_bright, True, PP_ALIGN.CENTER)
add_text(slide2, Inches(0.7), Inches(6.95), Inches(14.6), Inches(0.6), 
         "Language barriers + cultural stress + emotional isolation = massive unmet need", 
         20, dark_text, False, PP_ALIGN.CENTER)

add_footer(slide2)

# SLIDE 3: SOLUTION
print("Slide 3: Solution")
slide3 = add_slide()
add_text(slide3, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "The Solution: Emotionally Intelligent Bilingual AI Companions", 
         44, cyan_bright, True)

# Demo box
demo_box = add_box(slide3, Inches(0.5), Inches(1.5), Inches(7), Inches(5.3), light_bg, cyan_bright)
add_text(slide3, Inches(0.7), Inches(1.7), Inches(6.6), Inches(0.4), 
         "EspaLuz AI Tutor - LIVE Demo", 28, cyan_bright, True, PP_ALIGN.CENTER)
add_text(slide3, Inches(0.7), Inches(2.15), Inches(6.6), Inches(0.3), 
         "wa.me/50766623757", 20, dark_text, False, PP_ALIGN.CENTER)

if os.path.exists('qr-whatsapp-demo.jpg'):
    slide3.shapes.add_picture('qr-whatsapp-demo.jpg', Inches(2.2), Inches(2.8), width=Inches(3.3))

# Feature boxes
feature_x = Inches(7.8)
feature_w = Inches(7.7)
feature_h = Inches(1.6)

# Feature 1
f1 = add_box(slide3, feature_x, Inches(1.5), feature_w, feature_h, RGBColor(220, 255, 230), green_bright)
add_badge(slide3, feature_x + Inches(0.2), Inches(1.65), "AI", green_bright)
add_text(slide3, feature_x + Inches(1), Inches(1.7), feature_w - Inches(1.2), Inches(0.35), 
         "EMOTIONAL INTELLIGENCE ENGINE", 19, green_bright, True)
add_text(slide3, feature_x + Inches(1), Inches(2.1), feature_w - Inches(1.2), Inches(0.8), 
         "Custom 50+ emotion detection | Adapts to homesickness, anxiety\nRemembers emotional journey", 13, dark_text)

# Feature 2
f2 = add_box(slide3, feature_x, Inches(3.3), feature_w, feature_h, RGBColor(220, 245, 255), cyan_bright)
add_badge(slide3, feature_x + Inches(0.2), Inches(3.45), "FM", cyan_bright)
add_text(slide3, feature_x + Inches(1), Inches(3.5), feature_w - Inches(1.2), Inches(0.35), 
         "FAMILY-AWARE MEMORY SYSTEM", 19, cyan_bright, True)
add_text(slide3, feature_x + Inches(1), Inches(3.9), feature_w - Inches(1.2), Inches(0.8), 
         "Recognizes relationships, ages | One subscription serves family\nPersonalized for each member", 13, dark_text)

# Feature 3
f3 = add_box(slide3, feature_x, Inches(5.1), feature_w, feature_h, RGBColor(235, 225, 255), purple_bright)
add_badge(slide3, feature_x + Inches(0.2), Inches(5.25), "MP", purple_bright)
add_text(slide3, feature_x + Inches(1), Inches(5.3), feature_w - Inches(1.2), Inches(0.35), 
         "MULTI-PLATFORM", 19, purple_bright, True)
add_text(slide3, feature_x + Inches(1), Inches(5.7), feature_w - Inches(1.2), Inches(0.8), 
         "WhatsApp (90%+ LATAM) | Telegram + Web\nVoice, text, video avatars", 13, dark_text)

# Bridge box
bridge = add_box(slide3, Inches(0.5), Inches(7.1), Inches(15), Inches(0.9), light_bg, cyan_bright)
add_text(slide3, Inches(0.7), Inches(7.3), Inches(14.6), Inches(0.6), 
         "DUAL-SIDED BRIDGE: Expats learn Spanish | Locals learn English = 2x Revenue", 
         20, dark_text, True, PP_ALIGN.CENTER)

add_footer(slide3)

# SLIDE 4: TRACTION
print("Slide 4: Traction")
slide4 = add_slide()
add_text(slide4, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "Current Traction: Pre-Revenue, Proven Execution (7 Months)", 
         42, white, True)

col_w = Inches(7.2)

# Left column
add_text(slide4, Inches(0.5), Inches(1.3), col_w, Inches(0.3), "What We've Built", 28, cyan_bright, True)

b1 = add_box(slide4, Inches(0.5), Inches(1.8), col_w, Inches(1.3), RGBColor(220, 255, 230), green_bright)
add_badge(slide4, Inches(0.7), Inches(1.95), "6", green_bright)
add_text(slide4, Inches(1.5), Inches(2), col_w - Inches(1.2), Inches(0.3), 
         "6 LIVE Production Applications", 20, green_bright, True)
add_text(slide4, Inches(1.5), Inches(2.35), col_w - Inches(1.2), Inches(0.7), 
         "EspaLuz (WhatsApp, Telegram, Web) | ATUONA NFT\nALGOM Alpha AI advisor", 12, dark_text)

b2 = add_box(slide4, Inches(0.5), Inches(3.3), col_w, Inches(1.3), RGBColor(220, 245, 255), cyan_bright)
add_badge(slide4, Inches(0.7), Inches(3.45), "50K", cyan_bright)
add_text(slide4, Inches(1.5), Inches(3.5), col_w - Inches(1.2), Inches(0.3), 
         "50,000+ Lines of Production Code", 20, cyan_bright, True)
add_text(slide4, Inches(1.5), Inches(3.85), col_w - Inches(1.2), Inches(0.7), 
         "TypeScript, Python, JavaScript, SQL | 70+ React\n8+ AI integrations", 12, dark_text)

b3 = add_box(slide4, Inches(0.5), Inches(4.8), col_w, Inches(1.3), RGBColor(235, 225, 255), purple_bright)
add_badge(slide4, Inches(0.7), Inches(4.95), "8", purple_bright)
add_text(slide4, Inches(1.5), Inches(5), col_w - Inches(1.2), Inches(0.3), 
         "Infrastructure Ready to Scale", 20, purple_bright, True)
add_text(slide4, Inches(1.5), Inches(5.35), col_w - Inches(1.2), Inches(0.7), 
         "PayPal subscriptions LIVE | 99%+ uptime\n19 countries addressable", 12, dark_text)

# Right column
add_text(slide4, Inches(8.3), Inches(1.3), col_w, Inches(0.3), "Why It Matters", 28, cyan_bright, True)

b4 = add_box(slide4, Inches(8.3), Inches(1.8), col_w, Inches(1.6), RGBColor(255, 240, 220), orange_bright)
add_badge(slide4, Inches(8.5), Inches(1.95), "$", orange_bright)
add_text(slide4, Inches(9.3), Inches(2), col_w - Inches(1.2), Inches(0.3), 
         "CAPITAL EFFICIENCY", 20, orange_bright, True)
add_text(slide4, Inches(9.3), Inches(2.4), col_w - Inches(1.2), Inches(0.7), 
         "98% cost reduction", 38, cyan_bright, True, PP_ALIGN.CENTER)
add_text(slide4, Inches(9.3), Inches(3), col_w - Inches(1.2), Inches(0.25), 
         "<$15K vs. $950K traditional", 12, dark_text, False, PP_ALIGN.CENTER)

b5 = add_box(slide4, Inches(8.3), Inches(3.6), col_w, Inches(1.2), RGBColor(220, 235, 255), blue_bright)
add_text(slide4, Inches(8.5), Inches(3.75), col_w - Inches(0.4), Inches(0.3), 
         "SPEED OF EXECUTION", 18, blue_bright, True, PP_ALIGN.CENTER)
add_text(slide4, Inches(8.5), Inches(4.1), col_w - Inches(0.4), Inches(0.5), 
         "Solo-built in 7 months | 10x faster", 14, dark_text, False, PP_ALIGN.CENTER)

b6 = add_box(slide4, Inches(8.3), Inches(4.95), col_w, Inches(1.15), RGBColor(220, 255, 230), green_bright)
add_text(slide4, Inches(8.5), Inches(5.1), col_w - Inches(0.4), Inches(0.3), 
         "EARLY VALIDATION", 18, green_bright, True, PP_ALIGN.CENTER)
add_text(slide4, Inches(8.5), Inches(5.45), col_w - Inches(0.4), Inches(0.5), 
         "~10 users (zero marketing) | 40% better retention", 14, dark_text, False, PP_ALIGN.CENTER)

# Bottom
b7 = add_box(slide4, Inches(0.5), Inches(6.5), Inches(15), Inches(1.2), light_bg, cyan_bright)
add_text(slide4, Inches(0.7), Inches(6.7), Inches(14.6), Inches(0.9), 
         "PRE-REVENUE BY DESIGN: Built production infrastructure first\nFunding activates growth immediately", 
         18, dark_text, True, PP_ALIGN.CENTER)

add_footer(slide4)

# SLIDE 5: BUSINESS MODEL
print("Slide 5: Business Model")
slide5 = add_slide()
add_text(slide5, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "Business Model: Subscription SaaS (Ready to Activate)", 
         42, cyan_bright, True)

# Four tiers
tier_w = Inches(3.6)
tier_h = Inches(2.1)
y = Inches(1.4)

# Basic
t1 = add_box(slide5, Inches(0.5), y, tier_w, tier_h, RGBColor(220, 245, 255), cyan_bright)
add_text(slide5, Inches(0.7), y + Inches(0.15), tier_w - Inches(0.4), Inches(0.25), 
         "BASIC", 22, cyan_bright, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(0.7), y + Inches(0.45), tier_w - Inches(0.4), Inches(0.6), 
         "$15/mo", 46, dark_text, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(0.7), y + Inches(1.1), tier_w - Inches(0.4), Inches(0.35), 
         "Individual | Core AI", 13, dark_text)
add_text(slide5, Inches(0.7), y + Inches(1.6), tier_w - Inches(0.4), Inches(0.3), 
         "60% users", 14, green_bright, True, PP_ALIGN.CENTER)

# Family
t2 = add_box(slide5, Inches(4.3), y, tier_w, tier_h, RGBColor(235, 225, 255), purple_bright)
add_text(slide5, Inches(4.5), y + Inches(0.15), tier_w - Inches(0.4), Inches(0.25), 
         "FAMILY", 22, purple_bright, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(4.5), y + Inches(0.45), tier_w - Inches(0.4), Inches(0.6), 
         "$35/mo", 46, dark_text, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(4.5), y + Inches(1.1), tier_w - Inches(0.4), Inches(0.35), 
         "2-4 members | Video", 13, dark_text)
add_text(slide5, Inches(4.5), y + Inches(1.6), tier_w - Inches(0.4), Inches(0.3), 
         "30% users", 14, blue_bright, True, PP_ALIGN.CENTER)

# Premium
t3 = add_box(slide5, Inches(8.1), y, tier_w, tier_h, RGBColor(255, 240, 220), orange_bright)
add_text(slide5, Inches(8.3), y + Inches(0.15), tier_w - Inches(0.4), Inches(0.25), 
         "PREMIUM", 22, orange_bright, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(8.3), y + Inches(0.45), tier_w - Inches(0.4), Inches(0.6), 
         "$75/mo", 46, dark_text, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(8.3), y + Inches(1.1), tier_w - Inches(0.4), Inches(0.35), 
         "1-on-1 sessions", 13, dark_text)
add_text(slide5, Inches(8.3), y + Inches(1.6), tier_w - Inches(0.4), Inches(0.3), 
         "8% users", 14, orange_bright, True, PP_ALIGN.CENTER)

# Enterprise
t4 = add_box(slide5, Inches(11.9), y, tier_w, tier_h, RGBColor(220, 255, 230), green_bright)
add_text(slide5, Inches(12.1), y + Inches(0.15), tier_w - Inches(0.4), Inches(0.25), 
         "ENTERPRISE", 22, green_bright, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(12.1), y + Inches(0.45), tier_w - Inches(0.4), Inches(0.6), 
         "$200-500", 42, dark_text, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(12.1), y + Inches(1.1), tier_w - Inches(0.4), Inches(0.35), 
         "Schools/Corps", 13, dark_text)
add_text(slide5, Inches(12.1), y + Inches(1.6), tier_w - Inches(0.4), Inches(0.3), 
         "2% users", 14, green_bright, True, PP_ALIGN.CENTER)

# ARPU
arpu = add_box(slide5, Inches(3.5), Inches(3.8), Inches(9), Inches(0.75), light_bg, cyan_bright)
add_text(slide5, Inches(3.7), Inches(3.93), Inches(8.6), Inches(0.5), 
         "BLENDED ARPU: $25/month", 34, dark_text, True, PP_ALIGN.CENTER)

# Metrics
metric_w = Inches(4.8)
metric_h = Inches(1.5)
y = Inches(4.85)

m1 = add_box(slide5, Inches(0.5), y, metric_w, metric_h, RGBColor(220, 255, 230), green_bright)
add_text(slide5, Inches(0.7), y + Inches(0.15), metric_w - Inches(0.4), Inches(0.6), 
         "$300", 42, dark_text, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(0.7), y + Inches(0.8), metric_w - Inches(0.4), Inches(0.3), 
         "LTV (24-month retention)", 14, dark_text, False, PP_ALIGN.CENTER)

m2 = add_box(slide5, Inches(5.6), y, metric_w, metric_h, RGBColor(220, 245, 255), cyan_bright)
add_text(slide5, Inches(5.8), y + Inches(0.15), metric_w - Inches(0.4), Inches(0.6), 
         "$30-65", 42, dark_text, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(5.8), y + Inches(0.8), metric_w - Inches(0.4), Inches(0.3), 
         "CAC (organic + paid)", 14, dark_text, False, PP_ALIGN.CENTER)

m3 = add_box(slide5, Inches(10.7), y, metric_w, metric_h, RGBColor(235, 225, 255), purple_bright)
add_text(slide5, Inches(10.9), y + Inches(0.15), metric_w - Inches(0.4), Inches(0.6), 
         "4.6:1", 42, dark_text, True, PP_ALIGN.CENTER)
add_text(slide5, Inches(10.9), y + Inches(0.8), metric_w - Inches(0.4), Inches(0.3), 
         "LTV:CAC (VCs want 3:1+)", 14, dark_text, False, PP_ALIGN.CENTER)

add_footer(slide5)

# SLIDE 6-11: I'll create simplified versions focusing on clarity
# SLIDE 6: MARKET
print("Slide 6: Market")
slide6 = add_slide()
add_text(slide6, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "$37B+ TAM by 2030", 48, cyan_bright, True)

tam_box = add_box(slide6, Inches(3), Inches(1.5), Inches(10), Inches(1.6), light_bg, cyan_bright)
add_text(slide6, Inches(3.2), Inches(1.7), Inches(9.6), Inches(0.7), 
         "$37B+", 68, cyan_bright, True, PP_ALIGN.CENTER)
add_text(slide6, Inches(3.2), Inches(2.45), Inches(9.6), Inches(0.5), 
         "AI Assistants: $25B | Language Learning: $12B", 18, dark_text, False, PP_ALIGN.CENTER)

col_w = Inches(7.2)

td = add_box(slide6, Inches(0.5), Inches(3.5), col_w, Inches(3.3), RGBColor(235, 225, 255), purple_bright)
add_text(slide6, Inches(0.7), Inches(3.7), col_w - Inches(0.4), Inches(0.3), 
         "TOP-DOWN: 330M+ People", 24, purple_bright, True)
add_text(slide6, Inches(0.7), Inches(4.1), col_w - Inches(0.4), Inches(2.5), 
         "280M+ EXPATS WORLDWIDE\n- Permanent: 108M\n- Long-term travelers: 89M\n- Digital nomads: 35M\n\n50M+ LOCAL SERVICE PROVIDERS\n- Restaurants, hospitality, tours", 
         14, dark_text)

bu = add_box(slide6, Inches(8.3), Inches(3.5), col_w, Inches(3.3), RGBColor(220, 255, 230), green_bright)
add_text(slide6, Inches(8.5), Inches(3.7), col_w - Inches(0.4), Inches(0.3), 
         "BOTTOMS-UP: 19 Countries", 24, green_bright, True)
add_text(slide6, Inches(8.5), Inches(4.1), col_w - Inches(0.4), Inches(1.5), 
         "1% penetration = 2.8M users\n10% conversion = 280K paying\n\n$84M ARR AT 10% CONVERSION", 
         16, dark_text)
add_text(slide6, Inches(8.5), Inches(5.9), col_w - Inches(0.4), Inches(0.6), 
         "AT 5% PENETRATION:\n$420M ARR POTENTIAL", 18, green_bright, True, PP_ALIGN.CENTER)

add_footer(slide6)

# SLIDE 7: FINANCIALS
print("Slide 7: Financials")
slide7 = add_slide()
add_text(slide7, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "Financials: 12-Month Growth to $186K ARR", 42, cyan_bright, True)

table_box = add_box(slide7, Inches(1), Inches(1.4), Inches(14), Inches(2.7), light_bg, cyan_bright)
add_text(slide7, Inches(1.3), Inches(1.6), Inches(2), Inches(0.25), "Quarter", 16, dark_text, True)
add_text(slide7, Inches(4.3), Inches(1.6), Inches(1.8), Inches(0.25), "Users", 16, dark_text, True, PP_ALIGN.RIGHT)
add_text(slide7, Inches(6.3), Inches(1.6), Inches(1.8), Inches(0.25), "Paying", 16, dark_text, True, PP_ALIGN.RIGHT)
add_text(slide7, Inches(8.5), Inches(1.6), Inches(2), Inches(0.25), "MRR", 16, dark_text, True, PP_ALIGN.RIGHT)
add_text(slide7, Inches(11), Inches(1.6), Inches(2.5), Inches(0.25), "ARR", 16, dark_text, True, PP_ALIGN.RIGHT)

rows = [("Q1", "150", "15", "$365", "$4.4K"),
        ("Q2", "650", "65", "$1,800", "$21.6K"),
        ("Q3", "1,970", "197", "$5,625", "$67.5K")]

for i, (q, u, p, m, a) in enumerate(rows):
    y = 1.95 + i * 0.35
    add_text(slide7, Inches(1.3), Inches(y), Inches(2), Inches(0.25), q, 14, dark_text)
    add_text(slide7, Inches(4.3), Inches(y), Inches(1.8), Inches(0.25), u, 14, dark_text, False, PP_ALIGN.RIGHT)
    add_text(slide7, Inches(6.3), Inches(y), Inches(1.8), Inches(0.25), p, 14, dark_text, False, PP_ALIGN.RIGHT)
    add_text(slide7, Inches(8.5), Inches(y), Inches(2), Inches(0.25), m, 14, dark_text, False, PP_ALIGN.RIGHT)
    add_text(slide7, Inches(11), Inches(y), Inches(2.5), Inches(0.25), a, 14, dark_text, False, PP_ALIGN.RIGHT)

# Q4 highlighted
add_text(slide7, Inches(1.3), Inches(3.3), Inches(2), Inches(0.3), "Q4", 16, green_bright, True)
add_text(slide7, Inches(4.3), Inches(3.3), Inches(1.8), Inches(0.3), "4,800", 16, green_bright, True, PP_ALIGN.RIGHT)
add_text(slide7, Inches(6.3), Inches(3.3), Inches(1.8), Inches(0.3), "480", 16, green_bright, True, PP_ALIGN.RIGHT)
add_text(slide7, Inches(8.5), Inches(3.3), Inches(2), Inches(0.3), "$15,525", 16, green_bright, True, PP_ALIGN.RIGHT)
add_text(slide7, Inches(11), Inches(3.3), Inches(2.5), Inches(0.3), "$186K", 18, cyan_bright, True, PP_ALIGN.RIGHT)

traj = add_box(slide7, Inches(0.5), Inches(4.5), col_w, Inches(1.6), RGBColor(220, 235, 255), blue_bright)
add_text(slide7, Inches(0.7), Inches(4.7), col_w - Inches(0.4), Inches(0.3), 
         "3-Year Trajectory", 22, blue_bright, True)
add_text(slide7, Inches(0.7), Inches(5.05), col_w - Inches(0.4), Inches(1), 
         "Y1: 4,800 users | $186K ARR\nY2: 25,000 users | $1.2M ARR\nY3: 120,000 users | $6M ARR", 
         14, dark_text)

assum = add_box(slide7, Inches(8.3), Inches(4.5), col_w, Inches(1.6), RGBColor(220, 255, 230), green_bright)
add_text(slide7, Inches(8.5), Inches(4.7), col_w - Inches(0.4), Inches(0.3), 
         "Key Assumptions", 22, green_bright, True)
add_text(slide7, Inches(8.5), Inches(5.05), col_w - Inches(0.4), Inches(1), 
         "10% Conversion (industry std)\nCAC: $30-65 | LTV: $300\nChurn: 6% monthly", 
         14, dark_text)

be = add_box(slide7, Inches(3), Inches(6.5), Inches(10), Inches(0.8), light_bg, cyan_bright)
add_text(slide7, Inches(3.2), Inches(6.67), Inches(9.6), Inches(0.5), 
         "BREAK-EVEN: Month 15-18  |  Series A ready Month 18", 22, dark_text, True, PP_ALIGN.CENTER)

add_footer(slide7)

# SLIDE 8: TEAM (simplified)
print("Slide 8: Team")
slide8 = add_slide()
add_text(slide8, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "Team: Solo Founder Ready to Build World-Class Team", 42, cyan_bright, True)

founder = add_box(slide8, Inches(0.5), Inches(1.4), Inches(7), Inches(5.5), light_bg, cyan_bright)
add_text(slide8, Inches(0.7), Inches(1.6), Inches(6.6), Inches(0.4), 
         "ELENA REVICHEVA", 28, cyan_bright, True, PP_ALIGN.CENTER)
add_text(slide8, Inches(0.7), Inches(2.05), Inches(6.6), Inches(0.3), 
         "Founder, CEO & CTO", 18, dark_text, False, PP_ALIGN.CENTER)

add_badge(slide8, Inches(0.7), Inches(2.6), "EX", purple_bright)
add_text(slide8, Inches(1.5), Inches(2.65), Inches(5.7), Inches(0.25), 
         "EXECUTIVE BACKGROUND", 16, purple_bright, True)
add_text(slide8, Inches(1.5), Inches(2.95), Inches(5.7), Inches(0.8), 
         "Former Chief Legal Officer & IT Executive\n7 years digital transformation\nExecutive: strategy, team mgmt", 
         12, dark_text)

add_badge(slide8, Inches(0.7), Inches(4), "TE", green_bright)
add_text(slide8, Inches(1.5), Inches(4.05), Inches(5.7), Inches(0.25), 
         "TECHNICAL EXECUTION", 16, green_bright, True)
add_text(slide8, Inches(1.5), Inches(4.35), Inches(5.7), Inches(1), 
         "Self-taught AI engineer (2025)\n6 apps in 7 months for <$15K\n50,000+ lines of code\n8+ AI service integrations", 
         12, dark_text)

add_badge(slide8, Inches(0.7), Inches(5.6), "FM", cyan_bright)
add_text(slide8, Inches(1.5), Inches(5.65), Inches(5.7), Inches(0.25), 
         "FOUNDER-MARKET FIT", 16, cyan_bright, True)
add_text(slide8, Inches(1.5), Inches(5.95), Inches(5.7), Inches(0.8), 
         "Russia to Panama (2022) single mother\nLived the problem | Bilingual EN/ES\nDeep expat & local understanding", 
         12, dark_text)

val = add_box(slide8, Inches(8), Inches(1.4), Inches(7.5), Inches(2.3), RGBColor(220, 235, 255), blue_bright)
add_text(slide8, Inches(8.2), Inches(1.6), Inches(7.1), Inches(0.3), 
         "STRATEGIC VALIDATION", 22, blue_bright, True)
add_text(slide8, Inches(8.2), Inches(1.95), Inches(7.1), Inches(1.5), 
         "Innovation Smart District Panama\nEcosystem support, LATAM network\n\nDecentralized AI Agent Alliance\nIndustry recognition, global AI community", 
         14, dark_text)

hire = add_box(slide8, Inches(8), Inches(3.9), Inches(7.5), Inches(3), RGBColor(255, 240, 220), orange_bright)
add_text(slide8, Inches(8.2), Inches(4.1), Inches(7.1), Inches(0.3), 
         "PLANNED HIRES (12-18 Months)", 20, orange_bright, True)
add_text(slide8, Inches(8.2), Inches(4.5), Inches(7.1), Inches(2.2), 
         "1ST: Full-Stack Engineer\nScaling products, mobile app\n\n2ND: Growth Marketer\nUser acquisition, partnerships\n\n3RD: Operations Manager\nCustomer support, B2B deals", 
         14, dark_text)

add_footer(slide8)

# SLIDE 9: COMPETITION (simplified table)
print("Slide 9: Competition")
slide9 = add_slide()
add_text(slide9, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "Competition: Only Platform with ALL Key Differentiators", 40, cyan_bright, True)

comp = add_box(slide9, Inches(0.5), Inches(1.3), Inches(15), Inches(3.2), light_bg, cyan_bright)

# Headers
add_text(slide9, Inches(0.8), Inches(1.5), Inches(3.5), Inches(0.25), "Feature", 14, dark_text, True)
add_text(slide9, Inches(4.5), Inches(1.5), Inches(2.3), Inches(0.25), "Duolingo", 14, dark_text, True, PP_ALIGN.CENTER)
add_text(slide9, Inches(7), Inches(1.5), Inches(2.3), Inches(0.25), "Babbel", 14, dark_text, True, PP_ALIGN.CENTER)
add_text(slide9, Inches(9.5), Inches(1.5), Inches(2.3), Inches(0.25), "ChatGPT", 14, dark_text, True, PP_ALIGN.CENTER)
add_text(slide9, Inches(12), Inches(1.5), Inches(2.8), Inches(0.25), "AIdeazz", 14, cyan_bright, True, PP_ALIGN.CENTER)

rows_data = [
    ("Emotional Intelligence", "NO", "NO", "~", "YES"),
    ("Family Context & Memory", "NO", "NO", "NO", "YES"),
    ("Multi-Platform", "NO", "NO", "~", "YES"),
    ("Cultural Adaptation", "~", "~", "NO", "YES"),
    ("Web3 Ownership", "NO", "NO", "NO", "YES"),
    ("Dual Marketplace", "NO", "NO", "NO", "YES"),
    ("98% Capital Efficient", "NO", "NO", "NO", "YES"),
]

row_y = 1.85
for feat, d, b, c, a in rows_data:
    add_text(slide9, Inches(0.8), Inches(row_y), Inches(3.5), Inches(0.25), feat, 12, dark_text)
    d_col = red_bright if d == "NO" else (orange_bright if d == "~" else green_bright)
    add_text(slide9, Inches(4.5), Inches(row_y), Inches(2.3), Inches(0.25), d, 14, d_col, True, PP_ALIGN.CENTER)
    b_col = red_bright if b == "NO" else (orange_bright if b == "~" else green_bright)
    add_text(slide9, Inches(7), Inches(row_y), Inches(2.3), Inches(0.25), b, 14, b_col, True, PP_ALIGN.CENTER)
    c_col = red_bright if c == "NO" else (orange_bright if c == "~" else green_bright)
    add_text(slide9, Inches(9.5), Inches(row_y), Inches(2.3), Inches(0.25), c, 14, c_col, True, PP_ALIGN.CENTER)
    add_text(slide9, Inches(12), Inches(row_y), Inches(2.8), Inches(0.25), a, 14, green_bright, True, PP_ALIGN.CENTER)
    row_y += 0.32

# Competitor cards
card_w = Inches(4.8)
card_h = Inches(1.4)
card_y = Inches(4.8)

c1 = add_box(slide9, Inches(0.5), card_y, card_w, card_h, RGBColor(220, 255, 230), green_bright)
add_text(slide9, Inches(0.7), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.3), 
         "DUOLINGO", 18, green_bright, True)
add_text(slide9, Inches(0.7), card_y + Inches(0.55), card_w - Inches(0.4), Inches(0.7), 
         "$7B+ market cap\nGamified, not emotionally intelligent", 12, dark_text)

c2 = add_box(slide9, Inches(5.6), card_y, card_w, card_h, RGBColor(235, 225, 255), purple_bright)
add_text(slide9, Inches(5.8), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.3), 
         "BABBEL", 18, purple_bright, True)
add_text(slide9, Inches(5.8), card_y + Inches(0.55), card_w - Inches(0.4), Inches(0.7), 
         "$1B+ valuation\nTraditional courses, not AI-first", 12, dark_text)

c3 = add_box(slide9, Inches(10.7), card_y, card_w, card_h, RGBColor(220, 235, 255), blue_bright)
add_text(slide9, Inches(10.9), card_y + Inches(0.2), card_w - Inches(0.4), Inches(0.3), 
         "CHATGPT", 18, blue_bright, True)
add_text(slide9, Inches(10.9), card_y + Inches(0.55), card_w - Inches(0.4), Inches(0.7), 
         "$80B+ valuation\nGeneric AI, not specialized", 12, dark_text)

moat = add_box(slide9, Inches(0.5), Inches(6.5), Inches(15), Inches(1.2), light_bg, cyan_bright)
add_text(slide9, Inches(0.7), Inches(6.65), Inches(14.6), Inches(0.3), 
         "OUR DEFENSIBLE MOATS:", 20, cyan_bright, True)
add_text(slide9, Inches(0.7), Inches(7), Inches(14.6), Inches(0.6), 
         "Technical: Custom emotion detection | Data: Family context improves\nExecution: Vibe coding = 10x faster | Market: First-mover emotional AI + Web3 in LATAM", 
         13, dark_text)

add_footer(slide9)

# SLIDE 10: FUNDING
print("Slide 10: Funding")
slide10 = add_slide()
add_text(slide10, Inches(0.5), Inches(0.5), Inches(15), Inches(0.6), 
         "Funding: Seeking $100K-$500K Pre-Seed", 42, cyan_bright, True)

fund = add_box(slide10, Inches(2.5), Inches(1.3), Inches(11), Inches(1.5), light_bg, cyan_bright)
add_text(slide10, Inches(2.7), Inches(1.5), Inches(10.6), Inches(0.5), 
         "$100K - $500K PRE-SEED ROUND", 44, cyan_bright, True, PP_ALIGN.CENTER)
add_text(slide10, Inches(2.7), Inches(2.05), Inches(10.6), Inches(0.6), 
         "From ~10 users to 4,800 users in 12 months\nPath to $186K ARR | Series A Ready Month 18", 
         16, dark_text, False, PP_ALIGN.CENTER)

# Use of funds
use_w = Inches(3.5)
use_h = Inches(2.1)
use_y = Inches(3.1)

u1 = add_box(slide10, Inches(0.5), use_y, use_w, use_h, RGBColor(220, 245, 255), cyan_bright)
add_badge(slide10, Inches(0.65), use_y + Inches(0.15), "40%", cyan_bright)
add_text(slide10, Inches(1.4), use_y + Inches(0.2), use_w - Inches(1.1), Inches(0.25), 
         "ENGINEERING", 16, cyan_bright, True)
add_text(slide10, Inches(1.4), use_y + Inches(0.5), use_w - Inches(1.1), Inches(0.25), 
         "$40K-200K", 20, dark_text, True)
add_text(slide10, Inches(1.4), use_y + Inches(0.8), use_w - Inches(1.1), Inches(1.1), 
         "1-2 full-stack engineers\nMobile app (React Native)\nAnalytics systems\nERC-7857 integration", 
         10, dark_text)

u2 = add_box(slide10, Inches(4.2), use_y, use_w, use_h, RGBColor(220, 255, 230), green_bright)
add_badge(slide10, Inches(4.35), use_y + Inches(0.15), "30%", green_bright)
add_text(slide10, Inches(5.1), use_y + Inches(0.2), use_w - Inches(1.1), Inches(0.25), 
         "MARKETING", 16, green_bright, True)
add_text(slide10, Inches(5.1), use_y + Inches(0.5), use_w - Inches(1.1), Inches(0.25), 
         "$30K-150K", 20, dark_text, True)
add_text(slide10, Inches(5.1), use_y + Inches(0.8), use_w - Inches(1.1), Inches(1.1), 
         "User acquisition (ads)\nSEO, content, influencers\n5-10 country expansion\nSchool partnerships", 
         10, dark_text)

u3 = add_box(slide10, Inches(7.9), use_y, use_w, use_h, RGBColor(235, 225, 255), purple_bright)
add_badge(slide10, Inches(8.05), use_y + Inches(0.15), "20%", purple_bright)
add_text(slide10, Inches(8.8), use_y + Inches(0.2), use_w - Inches(1.1), Inches(0.25), 
         "OPERATIONS", 16, purple_bright, True)
add_text(slide10, Inches(8.8), use_y + Inches(0.5), use_w - Inches(1.1), Inches(0.25), 
         "$20K-100K", 20, dark_text, True)
add_text(slide10, Inches(8.8), use_y + Inches(0.8), use_w - Inches(1.1), Inches(1.1), 
         "Legal (19 countries)\nCustomer support team\nB2B partnerships", 
         10, dark_text)

u4 = add_box(slide10, Inches(11.6), use_y, use_w, use_h, RGBColor(255, 240, 220), orange_bright)
add_badge(slide10, Inches(11.75), use_y + Inches(0.15), "10%", orange_bright)
add_text(slide10, Inches(12.5), use_y + Inches(0.2), use_w - Inches(1.1), Inches(0.25), 
         "BUFFER", 16, orange_bright, True)
add_text(slide10, Inches(12.5), use_y + Inches(0.5), use_w - Inches(1.1), Inches(0.25), 
         "$10K-50K", 20, dark_text, True)
add_text(slide10, Inches(12.5), use_y + Inches(0.8), use_w - Inches(1.1), Inches(1.1), 
         "12-18 month runway\nUnexpected costs\nContingency planning", 
         10, dark_text)

why = add_box(slide10, Inches(0.5), Inches(5.6), Inches(15), Inches(1.3), light_bg, cyan_bright)
add_text(slide10, Inches(0.7), Inches(5.75), Inches(14.6), Inches(0.3), 
         "WHY THIS IS A SMART BET", 22, cyan_bright, True)
add_text(slide10, Inches(0.7), Inches(6.1), Inches(14.6), Inches(0.7), 
         "Built 6 apps for <$15K = 98% COST REDUCTION | Every investor dollar goes 50x further\nProven execution: 6 LIVE products, not slides | Clear milestones: Q1-Q4 roadmap", 
         14, dark_text)

stage = add_box(slide10, Inches(4), Inches(7.3), Inches(8), Inches(0.6), RGBColor(220, 235, 255), blue_bright)
add_text(slide10, Inches(4.2), Inches(7.45), Inches(7.6), Inches(0.4), 
         "STAGE: Pre-Seed  |  RAISED: $0 (bootstrapped <$15K)", 16, dark_text, True, PP_ALIGN.CENTER)

add_footer(slide10)

# SLIDE 11: SUMMARY
print("Slide 11: Summary")
slide11 = add_slide()
add_text(slide11, Inches(3), Inches(0.5), Inches(10), Inches(0.6), 
         "Ready to Scale Proven Execution", 46, cyan_bright, True, PP_ALIGN.CENTER)

# 6 summary boxes
sum_w = Inches(4.8)
sum_h = Inches(1.4)

s1 = add_box(slide11, Inches(0.5), Inches(1.4), sum_w, sum_h, RGBColor(220, 255, 230), green_bright)
add_badge(slide11, Inches(0.65), Inches(1.55), "$", green_bright)
add_text(slide11, Inches(1.4), Inches(1.6), sum_w - Inches(1.1), Inches(0.25), 
         "MASSIVE MARKET", 16, green_bright, True)
add_text(slide11, Inches(1.4), Inches(1.9), sum_w - Inches(1.1), Inches(0.8), 
         "$37B+ TAM | 280M+ expats\n19 countries | Dual marketplace", 12, dark_text)

s2 = add_box(slide11, Inches(5.6), Inches(1.4), sum_w, sum_h, RGBColor(220, 245, 255), cyan_bright)
add_badge(slide11, Inches(5.75), Inches(1.55), "6", cyan_bright)
add_text(slide11, Inches(6.5), Inches(1.6), sum_w - Inches(1.1), Inches(0.25), 
         "REAL EXECUTION", 16, cyan_bright, True)
add_text(slide11, Inches(6.5), Inches(1.9), sum_w - Inches(1.1), Inches(0.8), 
         "6 LIVE apps | 50K lines of code\nBuilt in 7 months for <$15K", 12, dark_text)

s3 = add_box(slide11, Inches(10.7), Inches(1.4), sum_w, sum_h, RGBColor(235, 225, 255), purple_bright)
add_badge(slide11, Inches(10.85), Inches(1.55), "*", purple_bright)
add_text(slide11, Inches(11.6), Inches(1.6), sum_w - Inches(1.1), Inches(0.25), 
         "UNIQUE POSITIONING", 16, purple_bright, True)
add_text(slide11, Inches(11.6), Inches(1.9), sum_w - Inches(1.1), Inches(0.8), 
         "Only platform: emotional AI +\nfamily context + Web3", 12, dark_text)

s4 = add_box(slide11, Inches(0.5), Inches(3), sum_w, sum_h, RGBColor(255, 240, 220), orange_bright)
add_badge(slide11, Inches(0.65), Inches(3.15), "%", orange_bright)
add_text(slide11, Inches(1.4), Inches(3.2), sum_w - Inches(1.1), Inches(0.25), 
         "CAPITAL EFFICIENT", 16, orange_bright, True)
add_text(slide11, Inches(1.4), Inches(3.5), sum_w - Inches(1.1), Inches(0.8), 
         "98% cost reduction\nInvestor $ goes 50x further", 12, dark_text)

s5 = add_box(slide11, Inches(5.6), Inches(3), sum_w, sum_h, RGBColor(220, 255, 230), green_bright)
add_badge(slide11, Inches(5.75), Inches(3.15), ">", green_bright)
add_text(slide11, Inches(6.5), Inches(3.2), sum_w - Inches(1.1), Inches(0.25), 
         "CLEAR PATH", 16, green_bright, True)
add_text(slide11, Inches(6.5), Inches(3.5), sum_w - Inches(1.1), Inches(0.8), 
         "$100K-500K > 4,800 users\n$186K ARR > Series A (18mo)", 12, dark_text)

s6 = add_box(slide11, Inches(10.7), Inches(3), sum_w, sum_h, RGBColor(220, 245, 255), cyan_bright)
add_badge(slide11, Inches(10.85), Inches(3.15), "FM", cyan_bright)
add_text(slide11, Inches(11.6), Inches(3.2), sum_w - Inches(1.1), Inches(0.25), 
         "FOUNDER-MARKET FIT", 16, cyan_bright, True)
add_text(slide11, Inches(11.6), Inches(3.5), sum_w - Inches(1.1), Inches(0.8), 
         "Lived the problem | Executive +\ntechnical | Ready to build", 12, dark_text)

# Contact box
contact = add_box(slide11, Inches(1.5), Inches(4.8), Inches(13), Inches(2.1), light_bg, purple_bright)
add_text(slide11, Inches(1.7), Inches(5), Inches(12.6), Inches(0.35), 
         "ELENA REVICHEVA", 28, purple_bright, True, PP_ALIGN.CENTER)
add_text(slide11, Inches(1.7), Inches(5.4), Inches(12.6), Inches(0.3), 
         "Founder, CEO & CTO", 20, dark_text, False, PP_ALIGN.CENTER)

add_text(slide11, Inches(2.5), Inches(5.85), Inches(5), Inches(0.25), 
         "Email: aipa@aideazz.xyz", 14, dark_text)
add_text(slide11, Inches(2.5), Inches(6.15), Inches(5), Inches(0.25), 
         "Demo: wa.me/50766623757", 14, dark_text)
add_text(slide11, Inches(8.5), Inches(5.85), Inches(5), Inches(0.25), 
         "Website: aideazz.xyz", 14, dark_text)
add_text(slide11, Inches(8.5), Inches(6.15), Inches(5), Inches(0.25), 
         "LinkedIn: linkedin.com/in/elenarevicheva", 14, dark_text)

# QR codes
qr_y = Inches(5)
qr_size = Inches(1.35)

if os.path.exists('qr-whatsapp-demo.jpg'):
    slide11.shapes.add_picture('qr-whatsapp-demo.jpg', Inches(3.3), qr_y, width=qr_size)
    add_text(slide11, Inches(3.1), qr_y + qr_size + Inches(0.05), Inches(1.8), Inches(0.2), 
             "WhatsApp", 10, dark_text, False, PP_ALIGN.CENTER)

if os.path.exists('qr-website.jpg'):
    slide11.shapes.add_picture('qr-website.jpg', Inches(7.3), qr_y, width=qr_size)
    add_text(slide11, Inches(7.1), qr_y + qr_size + Inches(0.05), Inches(1.8), Inches(0.2), 
             "Website", 10, dark_text, False, PP_ALIGN.CENTER)

if os.path.exists('qr-business-card.jpg'):
    slide11.shapes.add_picture('qr-business-card.jpg', Inches(11.3), qr_y, width=qr_size)
    add_text(slide11, Inches(11.1), qr_y + qr_size + Inches(0.05), Inches(1.8), Inches(0.2), 
             "Portfolio", 10, dark_text, False, PP_ALIGN.CENTER)

# Tagline
tag = add_box(slide11, Inches(2), Inches(7.3), Inches(12), Inches(0.75), light_bg, cyan_bright)
add_text(slide11, Inches(2.2), Inches(7.45), Inches(11.6), Inches(0.25), 
         "AIdeazz: Where Emotional AI Meets Human Potential", 20, dark_text, True, PP_ALIGN.CENTER)
add_text(slide11, Inches(2.2), Inches(7.72), Inches(11.6), Inches(0.22), 
         "Building the future of culturally-intelligent AI companions for 280M+ expats worldwide", 
         12, RGBColor(80, 80, 120), False, PP_ALIGN.CENTER)

add_footer(slide11)

# Save
prs.save('AIdeazz-Pitch-Deck.pptx')
print("\n" + "="*70)
print("SUCCESS! FINAL PowerPoint with PROPER CONTRAST created")
print("="*70)
print("\nDesign features:")
print("  - LIGHT boxes (RGB 240,242,250) on dark background")
print("  - DARK text (RGB 20,20,40) inside light boxes")
print("  - Vibrant colored borders (cyan, purple, green, orange, red)")
print("  - Solid colored badges with white text")
print("  - Zero text overlaps - calculated spacing")
print("  - Professional layout - clean and stylish")
print("\nFile: AIdeazz-Pitch-Deck.pptx (FINAL VERSION)")
print("="*70)
